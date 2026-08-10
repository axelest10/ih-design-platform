"""Renderizador determinista de presentaciones PPTX respaldadas por MaterialTemplate."""
from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from materials.models import MaterialTemplate, MaterialType

from .renderer import DEFAULT_LOGO_NAME, RenderValidationError
from .renderer_document import _approved_logo_path, _color_hex, _required_text


@dataclass(frozen=True)
class RenderedPresentation:
    template_key: str
    template_version: str
    data: dict[str, Any]
    asset_refs: list[str]
    validation_summary: dict[str, Any]
    pptx: bytes


def _rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor.from_string(value.upper())


def _add_text_box(slide, *, left, top, width, height, text, size, color, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Open Sans"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _build_pptx(*, headline: str, body: str, cta: str, logo_path, accent: str) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    white = _color_hex("white", default="white")
    dark = _color_hex("dark_navy", default="dark_navy")

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb(white)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.23), presentation.slide_height
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _rgb(accent)
    accent_bar.line.fill.background()

    slide.shapes.add_picture(str(logo_path), Inches(0.85), Inches(0.55), width=Inches(2.25))
    _add_text_box(
        slide,
        left=Inches(0.85),
        top=Inches(1.75),
        width=Inches(8.9),
        height=Inches(1.55),
        text=headline,
        size=50,
        color=accent,
        bold=True,
    )
    _add_text_box(
        slide,
        left=Inches(0.9),
        top=Inches(3.55),
        width=Inches(7.8),
        height=Inches(1.65),
        text=body,
        size=20,
        color=dark,
    )
    cta_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(5.65),
        Inches(3.2),
        Inches(0.62),
    )
    cta_shape.fill.solid()
    cta_shape.fill.fore_color.rgb = _rgb(accent)
    cta_shape.line.fill.background()
    cta_frame = cta_shape.text_frame
    cta_frame.clear()
    cta_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cta_run = cta_frame.paragraphs[0].add_run()
    cta_run.text = cta
    cta_run.font.name = "Open Sans"
    cta_run.font.size = Pt(15)
    cta_run.font.bold = True
    cta_run.font.color.rgb = _rgb(white)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def render_presentation_preview(
    payload: dict[str, Any], *, material_type: MaterialType
) -> RenderedPresentation:
    template_key = str(payload.get("template_key") or "presentation-16x9-v1")
    template = (
        MaterialTemplate.objects.filter(
            key=template_key,
            material_type=material_type,
            material_type__renderer_family=MaterialType.RendererFamily.PRESENTATION,
            active=True,
        )
        .select_related("material_type")
        .first()
    )
    if template is None:
        raise RenderValidationError(
            f"Template de presentación '{template_key}' no está disponible para este material."
        )
    if "pptx" not in template.output_formats:
        raise RenderValidationError(f"Template '{template_key}' no declara salida PPTX.")

    constraints = template.constraints or {}
    values = {
        field: _required_text(payload, field, constraints)
        for field in template.required_fields
    }
    headline = values.get("headline", str(payload.get("headline") or "").strip())
    body = values.get("body", str(payload.get("body") or "").strip())
    cta = values.get("cta", str(payload.get("cta") or "").strip())
    logo_name = str(payload.get("logo_name") or DEFAULT_LOGO_NAME)
    logo_path = _approved_logo_path(logo_name)
    accent_token = str(payload.get("accent_token") or "knowledge")
    accent_hex = _color_hex(accent_token, default="knowledge")
    pptx_bytes = _build_pptx(
        headline=headline,
        body=body,
        cta=cta,
        logo_path=logo_path,
        accent=accent_hex,
    )
    render_data = {
        "template_key": template.key,
        "template_version": template.version,
        "headline": headline,
        "body": body,
        "cta": cta,
        "logo_name": logo_name,
        "accent_token": accent_token,
    }
    validation_summary = {
        "status": "passed",
        "checks": [
            {"name": "template_registered", "status": "passed"},
            {"name": "required_fields", "status": "passed"},
            {"name": "logo_approved", "status": "passed", "asset": logo_name},
            {"name": "brand_color_authorized", "status": "passed", "hex": accent_hex},
            {"name": "pptx_generated", "status": "passed", "signature": "PK"},
        ],
    }
    return RenderedPresentation(
        template_key=template.key,
        template_version=template.version,
        data=render_data,
        asset_refs=[logo_name],
        validation_summary=validation_summary,
        pptx=pptx_bytes,
    )
