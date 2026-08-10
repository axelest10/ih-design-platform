from io import BytesIO

import pytest
from django.core.files.storage import FileSystemStorage
from pptx import Presentation
from rest_framework.test import APIClient

from briefs.models import DesignBrief
from designs.models import Design, DesignVersion
from designs.services.renderer import RenderValidationError
from designs.services.renderer_presentation import render_presentation_preview
from materials.models import MaterialType

PRESENTATION_PAYLOAD = {
    "template_key": "presentation-16x9-v1",
    "headline": "Inglés que conecta con el mundo",
    "body": "Una propuesta educativa internacional para cada etapa.",
    "cta": "Conoce nuestros programas",
    "logo_name": "ih-mexico-classic-png",
}


@pytest.mark.django_db
def test_presentation_renderer_generates_editable_pptx():
    material_type = MaterialType.objects.get(slug="presentation")

    rendered = render_presentation_preview(
        PRESENTATION_PAYLOAD,
        material_type=material_type,
    )

    assert rendered.pptx.startswith(b"PK")
    assert rendered.template_key == "presentation-16x9-v1"
    assert rendered.validation_summary["status"] == "passed"
    presentation = Presentation(BytesIO(rendered.pptx))
    assert len(presentation.slides) == 1
    assert "Inglés que conecta con el mundo" in " ".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )


@pytest.mark.django_db
def test_presentation_renderer_rejects_missing_required_text():
    material_type = MaterialType.objects.get(slug="presentation")
    incomplete = {**PRESENTATION_PAYLOAD}
    incomplete.pop("body")

    with pytest.raises(RenderValidationError, match="'body' es obligatorio"):
        render_presentation_preview(incomplete, material_type=material_type)


@pytest.mark.django_db
def test_presentation_preview_saves_pptx_with_configured_storage(tmp_path, monkeypatch):
    storage = FileSystemStorage(location=tmp_path, base_url="/media/")
    monkeypatch.setattr("designs.views.default_storage", storage)
    material_type = MaterialType.objects.get(slug="presentation")
    brief = DesignBrief.objects.create(
        title="Presentación piloto",
        format=DesignBrief.Format.HTML,
        material_type=material_type,
        audience="Escuelas",
        objective="Presentar la oferta educativa",
    )
    design = Design.objects.create(brief=brief)

    response = APIClient().post(
        f"/api/v1/designs/{design.pk}/preview/",
        PRESENTATION_PAYLOAD,
        format="json",
    )

    assert response.status_code == 201, response.json()
    assert response.json()["preview"]["pptx_url"].endswith(".pptx")
    version = DesignVersion.objects.get(design=design)
    pptx_path = version.render_data["pptx_path"]
    assert storage.exists(pptx_path)
    with storage.open(pptx_path, "rb") as pptx_file:
        assert pptx_file.read(2) == b"PK"
    assert pptx_path in version.asset_refs
